import docx
import filetype
import pptx
import pytesseract
import re
from PIL import Image
from openpyxl import load_workbook

def main(input):
    out = ""

    fileType = filetype.guess(input)
    if not fileType:
        out = open(input, encoding="utf-8", errors="ignore").read()
    else:

        print("\nDetected file type:", fileType.extension, "\n")

        if fileType.extension == "docx":
            document = docx.Document(input)
            for para in document.paragraphs:
                out += para.text + "\n"

        elif fileType.extension == "pptx":
            presentation = pptx.Presentation(input)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if not hasattr(shape, 'image'):
                        try:
                            out += re.sub('\n', ' ', shape.text) + "\n"
                        except:
                            pass

        elif fileType.extension in ["png", "pdf"]:
            out = pytesseract.image_to_string(Image.open(input))

        elif fileType.extension in ["md", "txt", "rtf"]:
            out = open(input, encoding="utf-8", errors="ignore").read()

        elif fileType.extension == "xlsx":
            wb = load_workbook(input, data_only=True)
            text_content = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text_content.append(f"\n--- Sheet: {sheet} ---\n")
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text_content.append(row_text)
            return "\n".join(text_content)

    # Clean up text a bit
    text = re.sub(r'[^\w\s.,!]+', ' ', out)
    print(text)
    return text
